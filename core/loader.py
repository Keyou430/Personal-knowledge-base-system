# -*- coding: utf-8 -*-
"""
文档加载器
支持 PDF、Word、PPT、Markdown、TXT、图片等多种格式
"""

import os
import re
import logging
from typing import List
from contextlib import contextmanager

from langchain_core.documents import Document
logger = logging.getLogger(__name__)


def load_pdf(file_path: str) -> List[Document]:
    """加载 PDF 文件，扫描件自动走 OCR"""
    from langchain_community.document_loaders import PyPDFLoader

    loader = PyPDFLoader(file_path)
    docs = loader.load()

    # 如果提取到的文字为空，说明是扫描件，走 OCR
    if docs and all(not d.page_content.strip() for d in docs):
        logger.info("PDF 文字提取为空，尝试 OCR 识别扫描件...")
        return _load_pdf_ocr(file_path)

    return docs


def _load_pdf_ocr(file_path: str) -> List[Document]:
    """
    将 PDF 每页渲染为图片，用 EasyOCR 提取文字
    """
    from core.image_processor import extract_text_ocr
    import pypdfium2 as pdfium
    from PIL import Image

    pdf = pdfium.PdfDocument(file_path)
    docs = []

    for i in range(len(pdf)):
        page = pdf[i]
        # 渲染为 PIL 图片（2x 分辨率提升 OCR 精度）
        bitmap = page.render(scale=2)
        image = bitmap.to_pil()

        # 保存为临时文件供 EasyOCR 使用
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
            image.save(tmp.name)
            tmp_path = tmp.name

        try:
            text = extract_text_ocr(tmp_path)
            if text.strip():
                docs.append(Document(
                    page_content=text,
                    metadata={"source": file_path, "page": i}
                ))
        finally:
            os.unlink(tmp_path)

    pdf.close()
    logger.info(f"PDF OCR 完成: {len(pdf)} 页 -> {len(docs)} 个文档片段")
    return docs


def load_docx(file_path: str) -> List[Document]:
    """Load DOCX blocks while preserving headings, tables, and source paths."""
    from docx import Document as DocxDocument

    doc = DocxDocument(file_path)
    blocks: List[Document] = []
    heading_path: list[str] = []
    table_index = 0

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = getattr(paragraph.style, "name", "") or ""
        if style_name.lower().startswith("heading"):
            try:
                level = int(style_name.split()[-1])
            except (TypeError, ValueError):
                level = 1
            heading_path = heading_path[: max(0, level - 1)]
            heading_path.append(text)
            blocks.append(
                Document(
                    page_content=text,
                    metadata={
                        "source": file_path,
                        "block_type": "heading",
                        "heading_level": level,
                        "heading_path": " / ".join(heading_path),
                    },
                )
            )
            continue
        blocks.append(
            Document(
                page_content=text,
                metadata={
                    "source": file_path,
                    "block_type": "paragraph",
                    "heading_path": " / ".join(heading_path),
                },
            )
        )

    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            blocks.append(
                Document(
                    page_content="\n".join(rows),
                    metadata={
                        "source": file_path,
                        "block_type": "table",
                        "table_index": table_index,
                        "heading_path": " / ".join(heading_path),
                    },
                )
            )
            table_index += 1

    return blocks


@contextmanager
def _com_app(prog_id: str):
    """
    COM 应用程序上下文管理器
    自动处理 CoInitialize/CoUninitialize 和应用进程的启停

    Args:
        prog_id: COM 程序标识符（如 "Word.Application"、"PowerPoint.Application"）
    """
    import win32com.client
    import pythoncom

    pythoncom.CoInitialize()
    app = None
    try:
        app = win32com.client.Dispatch(prog_id)
        app.Visible = False
        yield app
    finally:
        if app:
            try:
                app.Quit()
            except Exception:
                pass  # 进程可能已退出
        pythoncom.CoUninitialize()


def load_doc(file_path: str) -> List[Document]:
    """
    加载旧版 .doc 文件
    需要安装 Microsoft Word（通过 COM 自动化转换）
    """
    with _com_app("Word.Application") as word:
        doc = word.Documents.Open(os.path.abspath(file_path))
        try:
            text = doc.Content.Text
        finally:
            doc.Close(False)
    return [Document(page_content=text, metadata={"source": file_path})]


def load_pptx(file_path: str) -> List[Document]:
    """加载 .pptx 文件"""
    from pptx import Presentation

    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)

    full_text = "\n".join([t for t in texts if t.strip()])
    return [Document(page_content=full_text, metadata={"source": file_path})]


def load_ppt(file_path: str) -> List[Document]:
    """
    加载旧版 .ppt 文件
    需要安装 Microsoft PowerPoint（通过 COM 自动化转换）
    """
    with _com_app("PowerPoint.Application") as ppt:
        pres = ppt.Presentations.Open(os.path.abspath(file_path))
        try:
            texts = []
            for slide in pres.Slides:
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        texts.append(shape.TextFrame.TextRange.Text)
        finally:
            pres.Close()
    full_text = "\n".join([t for t in texts if t.strip()])
    return [Document(page_content=full_text, metadata={"source": file_path})]


def load_markdown(file_path: str) -> List[Document]:
    """Load Markdown with heading paths instead of flattening the file."""
    blocks: List[Document] = []
    heading_path: list[str] = []
    paragraph_lines: list[str] = []

    def flush() -> None:
        if paragraph_lines:
            text = "\n".join(paragraph_lines).strip()
            if text:
                blocks.append(
                    Document(
                        page_content=text,
                        metadata={
                            "source": file_path,
                            "block_type": "paragraph",
                            "heading_path": " / ".join(heading_path),
                        },
                    )
                )
            paragraph_lines.clear()

    with open(file_path, "r", encoding="utf-8") as handle:
        for raw_line in handle:
            line = raw_line.rstrip("\r\n")
            match = re.match(r"^(#{1,6})\s+(.+?)\s*$", line)
            if match:
                flush()
                level = len(match.group(1))
                heading_path = heading_path[: max(0, level - 1)]
                heading_path.append(match.group(2).strip())
                blocks.append(
                    Document(
                        page_content=match.group(2).strip(),
                        metadata={
                            "source": file_path,
                            "block_type": "heading",
                            "heading_level": level,
                            "heading_path": " / ".join(heading_path),
                        },
                    )
                )
            elif line.strip():
                paragraph_lines.append(line)
            else:
                flush()
    flush()
    return blocks


def load_text(file_path: str) -> List[Document]:
    """
    加载纯文本文件
    自动检测编码：先尝试 UTF-8，失败后回退到 GBK
    """
    for encoding in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            with open(file_path, "r", encoding=encoding) as handle:
                return [Document(page_content=handle.read(), metadata={"source": file_path})]
        except (UnicodeDecodeError, LookupError):
            continue
    # 所有编码都失败，使用 latin-1（永远不会报错，但可能乱码）
    with open(file_path, "r", encoding="latin-1") as handle:
        return [Document(page_content=handle.read(), metadata={"source": file_path})]


def load_image(file_path: str) -> List[Document]:
    """
    加载图片文件
    使用 EasyOCR 提取文字 + Qwen-VL 生成描述
    """
    from core.image_processor import process_image

    result = process_image(file_path)
    return [Document(
        page_content=result,
        metadata={"source": file_path, "type": "image"}
    )]


# 格式 -> 加载函数映射
_LOADER_MAP = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".doc": load_doc,
    ".pptx": load_pptx,
    ".ppt": load_ppt,
    ".md": load_markdown,
    ".txt": load_text,
    ".jpg": load_image,
    ".jpeg": load_image,
    ".png": load_image,
    ".bmp": load_image,
}


def load_document(file_path: str) -> List[Document]:
    """
    根据文件扩展名自动选择加载器

    Args:
        file_path: 文件路径

    Returns:
        Document 列表

    Raises:
        ValueError: 不支持的文件格式
        FileNotFoundError: 文件不存在
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()
    loader_func = _LOADER_MAP.get(ext)

    if loader_func is None:
        raise ValueError(f"不支持的文件格式: {ext}")

    logger.info(f"正在加载文件: {file_path} (格式: {ext})")
    docs = loader_func(file_path)
    logger.info(f"加载完成，共 {len(docs)} 个文档片段")
    return docs


def load_directory(dir_path: str) -> List[Document]:
    """
    批量加载目录下所有支持的文档

    Args:
        dir_path: 目录路径

    Returns:
        所有文档的 Document 列表
    """
    if not os.path.isdir(dir_path):
        raise NotADirectoryError(f"目录不存在: {dir_path}")

    all_docs = []
    for root, _, files in os.walk(dir_path):
        for filename in files:
            ext = os.path.splitext(filename)[1].lower()
            if ext in _LOADER_MAP:
                file_path = os.path.join(root, filename)
                try:
                    docs = load_document(file_path)
                    all_docs.extend(docs)
                except Exception as e:
                    logger.warning(f"加载文件失败: {file_path}, 错误: {e}")

    logger.info(f"目录加载完成，共 {len(all_docs)} 个文档片段")
    return all_docs
